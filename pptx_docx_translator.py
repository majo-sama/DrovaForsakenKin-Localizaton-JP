from llm_utils import get_llm_response
from typing import List, Dict
from collections import deque
from tqdm import tqdm
from pptx import Presentation
from docx import Document
from docx.text.paragraph import Paragraph
from textwrap import dedent


class TranslationService:
    MODEL_NAME = "gpt-4o-mini"

    @classmethod
    def sniff_writing_style(cls, sentences: List[str]) -> str:
        """Sniff the the sentences and return writing style."""
        samples = "\n".join(sentences)[:50]
        messages = [
            {
                "role": "user",
                "content": f"""
{samples}
---
Infer the scene in which the above sentences are likely used, and provide hints for translating the style and tone appropriately based on that.
The hints should be concise bullet points and within 100 characters.

Examples:

Shift\nLeft-Click\nRight-Click\nSPACE\nCtrl\nTAB\nENTER\nAlt\nUp\nDown\nLeft\nRight\nMove right\nMove left
=>
- **Context**: Likely a video game or software interface.
- **Tone**: Informal, engaging, and instructional.
- **Style**: Direct and concise; use action-oriented language.
- **Audience**: Gamers or tech users; assume familiarity with controls.
- **Pacing**: Quick, with emphasis on immediacy and responsiveness.
""",
            },
        ]
        response = get_llm_response(
            model_name=cls.MODEL_NAME,
            messages=messages,
            params_={"temperature": 0.0, "top_p": 0.3},
        )
        return response

    @classmethod
    def run(
        cls,
        sentences: List[str],
        source_lang: str = "English",
        target_lang: str = "Japanese",
        writing_style: str = "formal",
    ) -> List[str]:
        """Run translation on a list of sentences with context."""
        return cls.translate_with_context(
            sentences, source_lang, target_lang, writing_style
        )

    @classmethod
    def run_dict(
        cls,
        sentences: Dict[str, str],
        source_lang: str = "English",
        target_lang: str = "Japanese",
        writing_style: str = "formal",
    ) -> List[Dict[str, str]]:
        """Run translation on a dict of texts. The dict should have the following format:
        {"some_key": "text to translate", ...}

        The output will be in the same format.
        {"some_key": "translated text", ...}
        """
        translated_sentences = cls.run(
            list(sentences.values()),
            source_lang,
            target_lang,
            writing_style,
        )
        return {k: v for k, v in zip(sentences.keys(), translated_sentences)}

    @staticmethod
    def translate_sentence(
        sentence: str, previous_context: List[Dict[str, str]]
    ) -> str:
        """Translate a single sentence using a given context."""
        if not sentence.strip() or len(sentence.strip()) < 3:
            return sentence

        messages = previous_context + [{"role": "user", "content": sentence}]
        return get_llm_response(
            model_name=TranslationService.MODEL_NAME,
            messages=messages,
            params_={
                "temperature": 0.0,
                "top_p": 0.3,
                # "stop": [
                #     "<start_of_turn>",
                #     "<end_of_turn>",
                #     "\n",  # stop at the first newline because we are translating one sentence at a time
                # ],
            },
        )

    @staticmethod
    def translate_with_context(
        sentences: List[str],
        source_lang: str,
        target_lang: str,
        writing_style: str = "formal",
    ) -> List[str]:
        """Translate sentences considering context and writing style."""
        system_prompt = TranslationService.create_system_prompt(
            source_lang, target_lang
        )
        instruction_template = TranslationService.create_instruction_template(
            source_lang, target_lang, writing_style
        )

        instructions = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": instruction_template},
            {"role": "assistant", "content": "OK"},
        ]

        context = deque(maxlen=4)
        translated_sentences = []

        for sentence in tqdm(sentences):
            translation_context = instructions + list(context)
            translated_sentence = TranslationService.translate_sentence(
                sentence, translation_context
            )
            translated_sentences.append(translated_sentence)

            if sentence.strip():
                context.append({"role": "user", "content": sentence})
                context.append({"role": "assistant", "content": translated_sentence})

        return [s.strip() for s in translated_sentences]

    @staticmethod
    def create_system_prompt(source_lang: str, target_lang: str) -> str:
        """Create the system prompt for translation."""
        return dedent(f"""\
        You are a highly skilled professional {source_lang}-{target_lang} translator.
        Translate the given text accurately, considering the context and specific instructions provided.
        If no additional instructions or context are provided, use your expertise to provide a natural translation that aligns with the context.
        Strive to faithfully reflect the meaning and tone of the original text, pay attention to cultural nuances and ensure the translation is grammatically correct and easy to read.
        Review the translation for errors or unnatural expressions.
        For technical terms and proper nouns, use appropriate translations or leave them in the original language as necessary.""")

    @staticmethod
    def create_instruction_template(
        source_lang: str, target_lang: str, writing_style: str
    ) -> str:
        """Create the instruction template for the user."""
        return dedent(f"""\
        ## Instruction
        Translate {source_lang} into {target_lang}.
        Do not add any additional information or change the meaning of the text.

        Here are some input => output examples of English to Japanese translations:
        * Flash attention, with all these large language models and self-attention driving a transformer model, => フラッシュアテンション、最新の巨大言語モデルや自己注意機構がトランスフォーマーモデルを駆動させる中、
        * the idea behind flash attention is to dramatically accelerate the performance of these models  => フラッシュアテンションの背後にあるアイデアは、これらのモデルのパフォーマンスを劇的に高速化することです。
        * using an approach that I have personal experience with. => 私自身が経験した手法を用いて。

        When translating, please use the following hints: [writing_style:{writing_style}]
        Take a deep breath, calm down, and start translating.
        """)


class PptxTranslator:
    @staticmethod
    def run(
        input_pptx: str,
        output_pptx: str,
        source_lang: str = "Chinese",
        target_lang: str = "Japanese",
    ) -> List[Dict[str, str]]:
        """Run translation on a PPTX file and save the results."""
        text_boxes = PptxTranslator.extract_text_boxes(input_pptx)
        original_texts = [box["text"] for box in text_boxes]
        translated_texts = TranslationService.run(
            original_texts, source_lang, target_lang
        )

        PptxTranslator.replace_text_and_save(input_pptx, translated_texts, output_pptx)

        for i, box in enumerate(text_boxes):
            box["text"] = translated_texts[i]
        return text_boxes

    @staticmethod
    def extract_text_boxes(pptx_file: str) -> List[Dict[str, str]]:
        """Extract text boxes from a PPTX file."""
        presentation = Presentation(pptx_file)
        text_boxes = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_data = {
                        "text": shape.text,
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                        "page": slide.slide_id,
                    }
                    text_boxes.append(text_data)
        return text_boxes

    @staticmethod
    def replace_text_and_save(
        pptx_file: str, translated_texts: List[str], output_file: str
    ) -> None:
        """Replace text in a PPTX file with translated texts and save it."""
        presentation = Presentation(pptx_file)
        text_index = 0

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and text_index < len(translated_texts):
                    shape.text = translated_texts[text_index]
                    text_index += 1
        presentation.save(output_file)

    @staticmethod
    def save_as_markdown(boxes: List[Dict[str, str]], output_file: str) -> None:
        """Save the extracted text boxes as a Markdown file."""
        with open(output_file, "w") as f:
            current_page = None
            for box in boxes:
                if current_page != box["page"]:
                    f.write(f"\n\n### Page {box['page']}\n\n")
                    current_page = box["page"]
                f.write(f"{box['text']}\n")


class DocxTranslator:
    @staticmethod
    def run(
        input_docx: str,
        source_lang: str = "Chinese",
        target_lang: str = "Japanese",
    ) -> List[Paragraph]:
        """Run translation on a DOCX file and return translated paragraphs."""
        paragraphs = DocxTranslator.extract_paragraphs_from_docx(input_docx)
        original_texts = [para.text for para in paragraphs]
        translated_texts = TranslationService.run(
            original_texts, source_lang, target_lang
        )

        for i, para in enumerate(paragraphs):
            para.text = translated_texts[i]
        return paragraphs

    @staticmethod
    def extract_paragraphs_from_docx(docx_file: str) -> List[Paragraph]:
        """Extract paragraphs from a DOCX file."""
        document = Document(docx_file)
        return document.paragraphs

    @staticmethod
    def save_as_markdown(paragraphs: List[Paragraph], output_file: str) -> None:
        """Save paragraphs as a Markdown file."""
        with open(output_file, "w") as f:
            for para in paragraphs:
                f.write(f"{para.text}\n\n")


if __name__ == "__main__":
    # Example usage:

    # PPTX translation
    translated_boxes = PptxTranslator.run(
        "531159439894380757_2.4天恩師德.pptx",
        "531159439894380757_2.4天恩師德_ja.pptx",
        source_lang="Chinese",
        target_lang="Japanese",
    )
    PptxTranslator.save_as_markdown(
        translated_boxes, "531159439894380757_2.4天恩師德_ja.md"
    )

    # DOCX translation
    translated_paragraphs = DocxTranslator.run(
        "人生真諦.docx",
        source_lang="Chinese",
        target_lang="Japanese",
    )
    DocxTranslator.save_as_mark
